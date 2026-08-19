#!/usr/bin/env python3
"""Beanstech — NVIDIA Inception Pitch Deck (2026). Design: navy/gold, 16:9."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY2 = RGBColor(0x14, 0x2B, 0x4D)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
GOLD_L = RGBColor(0xF3, 0xE7, 0xC3)
CARD = RGBColor(0xF7, 0xF9, 0xFC)
BORDER = RGBColor(0xD9, 0xE1, 0xEC)
TEXT = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x5A, 0x6B, 0x7B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x0F, 0x5D, 0x50)

OUT = "/mnt/sda1/Parcerias/pldbr-tech/docs/Beanstech_Inception_Deck.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
prs.core_properties.title = "Beanstech — NVIDIA Inception Pitch Deck"
prs.core_properties.author = "BeansTech"


def slide_fundo(cor=NAVY):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = cor; bg.line.fill.background()
    return s


def caixa(s, x, y, w, h, fill=CARD, line=BORDER):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill; c.line.color.rgb = line; c.line.width = Pt(1)
    return c


def texto(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    first = True
    for txt, size, color, bold in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = align; first = False
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = "Arial"
    return tb


def bullets(s, x, y, w, h, itens, size=14, color=TEXT, gap=6):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(itens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            head, rest = it
            r = p.add_run(); r.text = "▸ " + head + " "
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Arial"
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = "Arial"
        else:
            r = p.add_run(); r.text = "▸ " + it
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Arial"
    return tb


def cabecalho(s, kicker, titulo):
    texto(s, 0.6, 0.35, 12.1, 0.4, [(kicker.upper(), 12, GOLD, True)])
    texto(s, 0.6, 0.62, 12.1, 0.9, [(titulo, 28, NAVY, True)])


# 1 — CAPA
s = slide_fundo(NAVY)
texto(s, 0.8, 2.2, 11.7, 0.5, [("BEANSTECH", 20, GOLD, True)])
texto(s, 0.8, 2.7, 11.7, 1.8, [("Trustworthy AI for Regulated Sectors", 40, WHITE, True)])
texto(s, 0.8, 4.4, 11.7, 0.9, [("Legal · Health · Finance — sovereign, citation-grounded intelligence", 16, GOLD_L, False)])
texto(s, 0.8, 6.5, 11.7, 0.5, [("NVIDIA Inception Program Application — 2026", 12, GOLD, True)])

# 2 — PROBLEMA
s = prs.slides.add_slide(BLANK)
cabecalho(s, "The problem", "In regulated sectors, wrong AI output is a legal liability")
cards = [
    ("Hallucination", "Generic LLMs invent case law, diagnoses and rules — unusable where citations are mandatory."),
    ("Data sovereignty", "Client data cannot flow through foreign clouds: LGPD, CFM 2.454, BACEN 3.978 constrain processing."),
    ("Unreadable sources", "Brazilian court systems produce millions of PDFs with corrupted fonts that no AI can read."),
]
for i, (t, d) in enumerate(cards):
    caixa(s, 0.6 + i * 4.1, 1.8, 3.8, 2.6)
    texto(s, 0.85 + i * 4.1, 2.0, 3.3, 0.6, [(t, 17, NAVY, True)])
    texto(s, 0.85 + i * 4.1, 2.6, 3.3, 1.7, [(d, 13, MUTED, False)])
texto(s, 0.6, 5.0, 12.1, 1.2, [("Result: law firms, hospitals and financial institutions in Brazil cannot adopt mainstream AI — ", 15, TEXT, False), ("the sectors that need trustworthy AI the most are the ones served the least.", 15, NAVY, True)])

# 3 — SOLUÇÃO
s = prs.slides.add_slide(BLANK)
cabecalho(s, "The solution", "One sovereign AI platform, three regulated verticals")
cards = [
    ("Legal AI", "RAG over 65M court decisions with citation-grounded answers and INPI-registered anti-hallucination protocol. Platforms: RAGJur, Advogando.AI, Minuta.Tech."),
    ("Health AI", "7 CFM-compliant platforms (FHIR R4): clinical AI scribe, exam support, digital prescription — compliance for the Aug/2026 CFM deadline."),
    ("RegTech", "Production AML engine — 25 detection typologies (BACEN 3.978) serving banks and public institutions."),
]
for i, (t, d) in enumerate(cards):
    caixa(s, 0.6 + i * 4.1, 1.8, 3.8, 3.4, fill=GOLD_L if i == 1 else CARD)
    texto(s, 0.85 + i * 4.1, 2.0, 3.3, 0.6, [(t, 17, NAVY, True)])
    texto(s, 0.85 + i * 4.1, 2.6, 3.3, 2.4, [(d, 12.5, TEXT, False)])
texto(s, 0.6, 5.6, 12.1, 0.8, [("Every answer cites its source. Every byte stays in Brazil. Zero telemetry by design.", 16, GREEN, True)])

# 4 — TRAÇÃO / MÉTRICAS
s = prs.slides.add_slide(BLANK)
cabecalho(s, "Traction", "Production systems, institutional validation")
metrics = [("65M+", "court decisions indexed"), ("55", "tribunals covered daily"), ("15+", "live platforms"), ("160", "specialized AI agents")]
for i, (n, l) in enumerate(metrics):
    caixa(s, 0.6 + i * 3.1, 1.8, 2.8, 1.8, fill=NAVY, line=NAVY)
    texto(s, 0.6 + i * 3.1, 2.0, 2.8, 0.8, [(n, 30, GOLD, True)], align=PP_ALIGN.CENTER)
    texto(s, 0.6 + i * 3.1, 2.85, 2.8, 0.6, [(l, 12, WHITE, False)], align=PP_ALIGN.CENTER)
bullets(s, 0.7, 4.1, 12.0, 2.6, [
    ("2 large clients in production", "— legal and financial institutions."),
    ("Hospital Albert Einstein", "— innovation team formally evaluating our architecture (their initiative)."),
    ("Google for Startups 2025 · Anthropic Partner", "— ecosystem recognitions; Oracle Startup program application submitted."),
    ("Public-sector offer", "— technology-transfer proposal to Brazil's Federal Prosecution Office under formal review."),
])

# 5 — DIFERENCIAIS
s = prs.slides.add_slide(BLANK)
cabecalho(s, "Differentiators", "Why institutions trust the stack")
bullets(s, 0.7, 1.9, 12.0, 4.8, [
    ("Citation-verified outputs —", "INPI-registered anti-hallucination protocol; every legal claim is checked against the source decision (0 false citations tolerated)."),
    ("Data sovereignty —", "on-premise or Brazil-region inference, zero telemetry: no token, metric or document ever leaves the institution's network."),
    ("Compliance by design —", "CFM 2.454 (health), BACEN 3.978 (AML), LGPD mapped per vertical before feature one."),
    ("Ingestion that reads Brazil —", "OCR engine that decodes corrupted court PDFs (PJe/e-SAJ); a 158-page case file is digitized in seconds."),
    ("Founder-market fit —", "lawyer (OAB) + Google Cloud Digital Leader building for the sectors he practiced in."),
], size=14, gap=10)

# 6 — NVIDIA TECHNOLOGY
s = prs.slides.add_slide(BLANK)
cabecalho(s, "NVIDIA technology", "Where CUDA runs our stack today — and tomorrow")
caixa(s, 0.6, 1.8, 5.9, 4.6)
texto(s, 0.85, 2.0, 5.4, 0.5, [("TODAY (production)", 15, NAVY, True)])
bullets(s, 0.85, 2.5, 5.4, 3.6, [
    "Open-weight LLM inference (Qwen/GLM, quantized) on NVIDIA L4/A100 (GCP, Brazil region)",
    "CUDA-accelerated embeddings + reranking over tens of millions of judgments",
    "Multimodal pipelines (document/image) on GPU workers",
], size=13)
caixa(s, 6.8, 1.8, 5.9, 4.6, fill=GOLD_L)
texto(s, 7.05, 2.0, 5.4, 0.5, [("ROADMAP WITH INCEPTION", 15, NAVY, True)])
bullets(s, 7.05, 2.5, 5.4, 3.6, [
    "TensorRT-LLM + vLLM: 2-4x latency/cost gains for legal drafting workloads",
    "Triton Inference Server for multi-tenant institutional serving",
    "On-prem GPU appliances (zero-telemetry deployments at courts, hospitals, banks)",
    "Open-weight video generation (Wan) for institutional training content",
], size=13)

# 7 — MERCADO
s = prs.slides.add_slide(BLANK)
cabecalho(s, "Market", "Brazil's regulated sectors, underserved by global AI")
bullets(s, 0.7, 1.9, 12.0, 4.4, [
    ("Legal —", "1.3M+ lawyers, 100M+ new case pages/year; citation duty makes generic LLMs legally unusable."),
    ("Health —", "CFM resolution 2.454 deadline (Aug/2026) forces every physician onto compliant tooling."),
    ("Finance —", "BACEN Circular 3.978 obliges AML typology coverage; mid-size institutions lack in-house AI."),
    (" wedge —", "we index what others can't read (court PDFs) and verify what others invent (citations)."),
], size=15, gap=12)

# 8 — TIME
s = prs.slides.add_slide(BLANK)
cabecalho(s, "Team", "Regulated-domain depth + production engineering")
caixa(s, 0.6, 1.8, 5.9, 4.2)
texto(s, 0.85, 2.0, 5.4, 0.5, [("Matheus Feijão — Founder & CEO", 15, NAVY, True)])
bullets(s, 0.85, 2.55, 5.4, 3.2, [
    "Lawyer (OAB/DF) — 13 years in federal justice system",
    "Google Cloud Digital Leader · Google for Startups founder",
    "Anthropic Partner · GLM/Qwen production deployments",
    "M.Sc. candidate (Digital Law, IDP-SP)",
], size=12.5)
caixa(s, 6.8, 1.8, 5.9, 4.2)
texto(s, 7.05, 2.0, 5.4, 0.5, [("Organization", 15, NAVY, True)])
bullets(s, 7.05, 2.55, 5.4, 3.2, [
    "5 people — 2 technical (engineering + security)",
    "160 internal AI agents automating product development",
    "GCP-based infra (Brazil region) with on-prem delivery muscle",
], size=12.5)

# 9 — ROADMAP + ASK
s = prs.slides.add_slide(BLANK)
cabecalho(s, "Roadmap & ask", "What Inception unlocks")
bullets(s, 0.7, 1.9, 12.0, 2.4, [
    ("2026 H2 —", "on-premise inference GA (TensorRT-LLM); 3 institutional pilots (justice, health, finance)."),
    ("2027 —", "multi-tenant sovereign AI serving; marketplace compliance modules; 25+ institutions."),
], size=15, gap=10)
caixa(s, 0.6, 4.3, 12.1, 2.3, fill=NAVY, line=NAVY)
texto(s, 0.9, 4.5, 11.5, 1.9, [
    ("We ask Inception for:", 15, GOLD, True),
    ("GPU discounts and DGX Cloud access for on-prem appliance certification · TensorRT-LLM/Triton engineering support · DLI training for the technical team · co-marketing to Brazil's regulated institutions.", 14, WHITE, False),
])

# 10 — CONTATO
s = slide_fundo(NAVY)
texto(s, 0.8, 2.6, 11.7, 1.2, [("Let's make regulated AI trustworthy — together.", 32, WHITE, True)])
texto(s, 0.8, 4.0, 11.7, 1.4, [
    ("Beanstech — Beans Tech Inova Simples (I.S.) · São Paulo, Brazil", 15, GOLD_L, False),
    ("matheus@beanstech.com.br · beanstech.com.br · linkedin.com/in/mfeijao · (11) 96650-7100", 14, GOLD, False),
])

prs.save(OUT)
print("OK:", OUT, "—", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
